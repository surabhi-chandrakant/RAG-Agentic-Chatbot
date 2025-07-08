import os
import json
import uuid
import asyncio
from typing import Dict, List, Any, Optional
from dataclasses import dataclass, asdict
from datetime import datetime
import logging
from pathlib import Path

# Document processing imports
import PyPDF2
import pandas as pd
from docx import Document
from pptx import Presentation
import faiss
import numpy as np
from sentence_transformers import SentenceTransformer

# Web framework
import streamlit as st
from streamlit.runtime.uploaded_file_manager import UploadedFile

# LLM Integration (Placeholder for actual LLM client)
try:
    import google.generativeai as genai
    # Configure your API key
    # genai.configure(api_key=os.environ.get("GOOGLE_API_KEY")) # Recommended: Load from environment variable
    LLM_AVAILABLE = True
except ImportError:
    LLM_AVAILABLE = False
    logging.warning("Google Generative AI library not found. LLM responses will be simulated.")
except Exception as e:
    LLM_AVAILABLE = False
    logging.error(f"Error configuring Google Generative AI: {e}. LLM responses will be simulated.")

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# --- MCP Message Protocol Classes (as provided in your code) ---
@dataclass
class MCPMessage:
    """Model Context Protocol message structure"""
    sender: str
    receiver: str
    type: str
    trace_id: str
    payload: Dict[str, Any]
    timestamp: str = None
    
    def __post_init__(self):
        if self.timestamp is None:
            self.timestamp = datetime.now().isoformat()
            
    def to_dict(self) -> Dict[str, Any]:
        return asdict(self)

class MCPMessageBus:
    """In-memory message bus for MCP communication"""
    def __init__(self):
        self.messages: List[MCPMessage] = []
        self.subscribers: Dict[str, List[callable]] = {}
        self.mcp_log_messages: List[Dict] = []

    def subscribe(self, agent_name: str, callback: callable):
        """Subscribe an agent to receive messages"""
        if agent_name not in self.subscribers:
            self.subscribers[agent_name] = []
        self.subscribers[agent_name].append(callback)
    
    async def publish(self, message: MCPMessage):
        """Publish a message to the bus and log for UI"""
        self.messages.append(message)
        logger.info(f"MCP Message: {message.sender} -> {message.receiver} ({message.type}) | Trace ID: {message.trace_id}")
        
        # Add to MCP log for UI display
        self.mcp_log_messages.append(message.to_dict())
        MAX_MCP_LOG_MESSAGES = 20
        if len(self.mcp_log_messages) > MAX_MCP_LOG_MESSAGES:
            self.mcp_log_messages = self.mcp_log_messages[-MAX_MCP_LOG_MESSAGES:]
        
        # Deliver to subscribers
        if message.receiver in self.subscribers:
            for callback in self.subscribers[message.receiver]:
                await callback(message)
    
    def get_trace_messages(self, trace_id: str) -> List[MCPMessage]:
        """Get all messages for a specific trace"""
        return [msg for msg in self.messages if msg.trace_id == trace_id]

    def get_mcp_log(self) -> List[Dict]:
        """Get the messages for MCP log display"""
        return self.mcp_log_messages

# --- Document Processing Classes (as provided in your code) ---
class DocumentProcessor:
    """Base class for document processing"""
    
    @staticmethod
    def process_pdf(file_content: bytes) -> str:
        """Extract text from PDF"""
        try:
            import io
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            return ""
    
    @staticmethod
    def process_docx(file_content: bytes) -> str:
        """Extract text from DOCX"""
        try:
            import io
            doc = Document(io.BytesIO(file_content))
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            return ""
    
    @staticmethod
    def process_pptx(file_content: bytes) -> str:
        """Extract text from PPTX"""
        try:
            import io
            prs = Presentation(io.BytesIO(file_content))
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return ""
    
    @staticmethod
    def process_csv(file_content: bytes) -> str:
        """Extract text from CSV"""
        try:
            import io
            df = pd.read_csv(io.BytesIO(file_content))
            text = f"CSV Data Summary:\n"
            text += f"Columns: {', '.join(df.columns.tolist())}\n"
            text += f"Rows: {len(df)}\n\n"
            text += "Sample Data:\n"
            text += df.head(10).to_string()
            text += "\n\nStatistical Summary:\n"
            text += df.describe().to_string()
            return text
        except Exception as e:
            logger.error(f"Error processing CSV: {e}")
            return ""
    
    @staticmethod
    def process_txt(file_content: bytes) -> str:
        """Extract text from TXT/Markdown"""
        try:
            return file_content.decode('utf-8')
        except Exception as e:
            logger.error(f"Error processing TXT: {e}")
            return ""

# --- Agent Classes (as provided in your code, with minor corrections/additions) ---
class IngestionAgent:
    """Agent responsible for document parsing and preprocessing"""
    
    def __init__(self, message_bus: MCPMessageBus):
        self.message_bus = message_bus
        self.message_bus.subscribe("IngestionAgent", self.handle_message)
        self.processor = DocumentProcessor()
        
    async def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "INGEST_DOCUMENT":
            await self.process_document(message)
    
    async def process_document(self, message: MCPMessage):
        """Process uploaded document"""
        filename = message.payload.get("filename", "unknown_file")
        try:
            file_info = message.payload
            file_content = file_info.get("content")
            file_type = file_info.get("type")
            
            text = ""
            if file_type == "application/pdf":
                text = self.processor.process_pdf(file_content)
            elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                text = self.processor.process_docx(file_content)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                text = self.processor.process_pptx(file_content)
            elif file_type == "text/csv":
                text = self.processor.process_csv(file_content)
            elif file_type in ["text/plain", "text/markdown"]:
                text = self.processor.process_txt(file_content)
            else:
                text = "Unsupported file format"

            logger.info(f"IngestionAgent: Processed {filename}. Extracted text length: {len(text)}")
            
            if not text or len(text.strip()) < 50:
                logger.warning(f"IngestionAgent: Extracted text from {filename} appears to be empty or very short")
                error_message = MCPMessage(
                    sender="IngestionAgent",
                    receiver="CoordinatorAgent",
                    type="ERROR",
                    trace_id=message.trace_id,
                    payload={"error": f"Failed to extract meaningful text from {filename}"}
                )
                await self.message_bus.publish(error_message)
                return

            chunks = self.chunk_text(text)
            logger.info(f"IngestionAgent: Generated {len(chunks)} chunks for {filename}")
            
            if not chunks:
                logger.warning(f"IngestionAgent: No chunks generated for {filename}")
                error_message = MCPMessage(
                    sender="IngestionAgent",
                    receiver="CoordinatorAgent",
                    type="ERROR",
                    trace_id=message.trace_id,
                    payload={"error": f"No chunks generated from {filename}"}
                )
                await self.message_bus.publish(error_message)
                return

            response_message = MCPMessage(
                sender="IngestionAgent",
                receiver="RetrievalAgent",
                type="DOCUMENT_PROCESSED",
                trace_id=message.trace_id,
                payload={
                    "filename": filename,
                    "chunks": chunks,
                    "original_text_preview": text[:1000] + "..." if len(text) > 1000 else text
                }
            )
            
            await self.message_bus.publish(response_message)
            
        except Exception as e:
            logger.error(f"Error in IngestionAgent processing document {filename}: {e}", exc_info=True)
            error_message = MCPMessage(
                sender="IngestionAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=message.trace_id,
                payload={"error": f"Error processing {filename}: {str(e)}"}
            )
            await self.message_bus.publish(error_message)
    
    def chunk_text(self, text: str, chunk_size: int = 250, overlap: int = 50) -> List[str]:
        """Split text into chunks for embedding with overlap."""
        if not text:
            return []
            
        words = text.split()
        chunks = []
        i = 0
        while i < len(words):
            chunk_end = min(i + chunk_size, len(words))
            chunk = " ".join(words[i:chunk_end])
            chunks.append(chunk)
            i += (chunk_size - overlap)
            if i >= len(words):
                break
            if chunk_end == len(words):
                break
            
        if not chunks and text.strip(): # Ensure at least one chunk if text exists
            chunks.append(text.strip())
            
        return chunks

class RetrievalAgent:
    """Agent responsible for embeddings and semantic retrieval"""
    
    def __init__(self, message_bus: MCPMessageBus):
        self.message_bus = message_bus
        self.message_bus.subscribe("RetrievalAgent", self.handle_message)
        # Explicitly set device to 'cpu' to avoid potential GPU memory issues
        self.embedder = SentenceTransformer('all-MiniLM-L6-v2', device='cpu')
        self.index = None
        self.chunks = [] # Stores original text chunks
        self.chunk_metadata = [] # Stores metadata including filename for each chunk
        
    async def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "DOCUMENT_PROCESSED":
            await self.embed_document(message)
        elif message.type == "RETRIEVE_CONTEXT":
            await self.retrieve_context(message)
        elif message.type == "CLEAR_INDEX":
            self.clear_index()
            await self.message_bus.publish(MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="INDEX_CLEARED",
                trace_id=message.trace_id,
                payload={"message": "RetrievalAgent index and data cleared successfully."}
            ))
            
    def clear_index(self):
        """Clears the FAISS index and associated data."""
        self.index = None
        self.chunks = []
        self.chunk_metadata = []
        logger.info("RetrievalAgent: FAISS index and all document data cleared.")

    async def embed_document(self, message: MCPMessage):
        """Create embeddings for document chunks"""
        filename = message.payload.get("filename", "unknown_file")
        try:
            chunks_to_embed = message.payload.get("chunks", [])

            if not chunks_to_embed:
                logger.warning(f"RetrievalAgent: Received no chunks for {filename}")
                error_message = MCPMessage(
                    sender="RetrievalAgent",
                    receiver="CoordinatorAgent",
                    type="ERROR",
                    trace_id=message.trace_id,
                    payload={"error": f"No chunks provided for {filename}"}
                )
                await self.message_bus.publish(error_message)
                return

            logger.info(f"RetrievalAgent: Embedding {len(chunks_to_embed)} chunks for {filename}")
            embeddings = self.embedder.encode(chunks_to_embed, convert_to_numpy=True)
            
            if self.index is None:
                dimension = embeddings.shape[1]
                self.index = faiss.IndexFlatL2(dimension)
                logger.info(f"RetrievalAgent: Initialized FAISS index with dimension {dimension}")
            
            start_idx = len(self.chunks) # Keep track of where new chunks start
            self.chunks.extend(chunks_to_embed)

            for i, chunk_text in enumerate(chunks_to_embed):
                self.chunk_metadata.append({
                    "filename": filename,
                    "chunk_index": start_idx + i, # Global index in self.chunks
                    "text": chunk_text
                })
            
            self.index.add(embeddings.astype('float32'))
            logger.info(f"RetrievalAgent: Added {len(embeddings)} embeddings to FAISS index. Total: {self.index.ntotal}")
            
            success_message = MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="EMBEDDING_COMPLETE",
                trace_id=message.trace_id,
                payload={
                    "filename": filename,
                    "chunks_processed": len(chunks_to_embed),
                    "total_chunks_in_index": self.index.ntotal
                }
            )
            
            await self.message_bus.publish(success_message)
            
        except Exception as e:
            logger.error(f"Error in RetrievalAgent embedding document {filename}: {e}", exc_info=True)
            error_message = MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=message.trace_id,
                payload={"error": f"Failed to embed document {filename}: {str(e)}"}
            )
            await self.message_bus.publish(error_message)
    
    async def retrieve_context(self, message: MCPMessage):
        """Retrieve relevant context for a query"""
        query = message.payload.get("query", "")
        top_k = message.payload.get("top_k", 5)
        
        try:
            logger.info(f"RetrievalAgent: Retrieving context for query: '{query}'")

            if self.index is None or self.index.ntotal == 0:
                logger.warning("RetrievalAgent: No documents indexed yet")
                response_message = MCPMessage(
                    sender="RetrievalAgent",
                    receiver="LLMResponseAgent",
                    type="CONTEXT_RESPONSE",
                    trace_id=message.trace_id,
                    payload={
                        "top_chunks": [],
                        "query": query,
                        "message": "No documents have been indexed yet."
                    }
                )
                await self.message_bus.publish(response_message)
                return
            
            query_embedding = self.embedder.encode([query], convert_to_numpy=True)
            actual_top_k = min(top_k, self.index.ntotal)
            distances, indices = self.index.search(query_embedding.astype('float32'), actual_top_k)
            
            top_chunks = []
            for i, idx in enumerate(indices[0]):
                if 0 <= idx < len(self.chunk_metadata):
                    metadata = self.chunk_metadata[idx]
                    top_chunks.append({
                        "text": metadata["text"],
                        "filename": metadata["filename"],
                        "score": float(distances[0][i])
                    })

            logger.info(f"RetrievalAgent: Retrieved {len(top_chunks)} chunks for query")

            response_message = MCPMessage(
                sender="RetrievalAgent",
                receiver="LLMResponseAgent",
                type="CONTEXT_RESPONSE",
                trace_id=message.trace_id,
                payload={
                    "top_chunks": top_chunks,
                    "query": query
                }
            )
            
            await self.message_bus.publish(response_message)
            
        except Exception as e:
            logger.error(f"Error in RetrievalAgent retrieval: {e}", exc_info=True)
            error_message = MCPMessage(
                sender="RetrievalAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=message.trace_id,
                payload={"error": f"Failed to retrieve context: {str(e)}"}
            )
            await self.message_bus.publish(error_message)

class LLMResponseAgent:
    """Agent responsible for generating LLM responses"""
    
    def __init__(self, message_bus: MCPMessageBus):
        self.message_bus = message_bus
        self.message_bus.subscribe("LLMResponseAgent", self.handle_message)
        self.model = None
        if LLM_AVAILABLE:
            try:
                # Ensure you have GOOGLE_API_KEY set in your environment variables
                genai.configure(api_key=os.environ.get("GOOGLE_API_KEY"))
                self.model = genai.GenerativeModel('gemini-1.5-flash')# Or other suitable model
                logger.info("LLMResponseAgent: Google Generative AI model loaded.")
            except Exception as e:
                logger.error(f"Failed to load Google Generative AI model: {e}", exc_info=True)
                self.model = None
        else:
            logger.warning("LLMResponseAgent: LLM not available. Responses will be simulated.")
            
    async def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        if message.type == "CONTEXT_RESPONSE":
            await self.generate_response(message)
    
    async def generate_response(self, message: MCPMessage):
        """Generate response using retrieved context"""
        query = message.payload.get("query", "")
        trace_id = message.trace_id
        
        try:
            top_chunks = message.payload.get("top_chunks", [])
            
            logger.info(f"LLMResponseAgent: Generating response for query '{query}' with {len(top_chunks)} chunks")

            # Format context
            context = ""
            sources = []
            
            if top_chunks:
                for chunk in top_chunks:
                    if 'text' in chunk and 'filename' in chunk:
                        context += f"Document: {chunk['filename']}\nContent: {chunk['text']}\n\n"
                        sources.append(chunk['filename'])
            else:
                context = "No relevant context found in uploaded documents."
            
            # Generate response using LLM or fallback
            response_text = await self.generate_rag_response(query, context, top_chunks)
            
            final_message = MCPMessage(
                sender="LLMResponseAgent",
                receiver="CoordinatorAgent",
                type="FINAL_RESPONSE",
                trace_id=trace_id,
                payload={
                    "query": query,
                    "response": response_text,
                    "sources": list(set(sources)), # Unique sources
                    "context_chunks_count": len(top_chunks)
                }
            )
            
            await self.message_bus.publish(final_message)
            
        except Exception as e:
            logger.error(f"Error in LLMResponseAgent: {e}", exc_info=True)
            error_message = MCPMessage(
                sender="LLMResponseAgent",
                receiver="CoordinatorAgent",
                type="ERROR",
                trace_id=trace_id,
                payload={"error": f"Failed to generate response: {str(e)}"}
            )
            await self.message_bus.publish(error_message)
    
    async def generate_rag_response(self, query: str, context: str, chunks: List[Dict]) -> str:
        """
        Generate a comprehensive response using the retrieved context.
        This method integrates with an actual LLM.
        """
        if not chunks:
            return f"I don't have any relevant information in the indexed documents to answer your question about '{query}'. Please upload and process relevant documents first."
            
        if self.model is None:
            # Fallback to rule-based or simple concatenation if LLM is not available
            logger.warning("LLM not available, synthesizing answer with rule-based approach.")
            return self.synthesize_answer(query, chunks, context) # Use the old rule-based method

        prompt = f"""You are an intelligent assistant. Use the following pieces of retrieved context to answer the question. 
        If you don't know the answer, just say that you don't have enough information from the provided documents. 
        Do not make up an answer.

        Context:
        {context}

        Question: {query}

        Answer:
        """
        
        try:
            # Using asyncio.to_thread for blocking LLM calls in an async environment
            # If the LLM client supports async, use await self.model.generate_content(prompt) directly
            response = await asyncio.to_thread(self.model.generate_content, prompt)
            return response.text
        except Exception as e:
            logger.error(f"Error calling LLM for RAG response: {e}", exc_info=True)
            return f"An error occurred while trying to generate a response. Please try again. (LLM Error: {str(e)})"
            
    def synthesize_answer(self, query: str, chunks: List[Dict], combined_text: str) -> str:
        """
        Synthesize a direct answer from the retrieved chunks.
        This is a rule-based approach that can be enhanced with actual LLM integration.
        Used as a fallback if LLM is not available.
        """
        query_lower = query.lower()
        
        # Try to provide a direct answer based on query patterns
        if "requirements" in query_lower or "requirement" in query_lower:
            return self.extract_requirements(chunks, combined_text)
        elif "citation" in query_lower or "cite" in query_lower:
            return self.extract_citation_info(chunks, combined_text)
        elif "what is" in query_lower or "define" in query_lower:
            return self.extract_definition(query, chunks, combined_text)
        elif "how to" in query_lower or "how do" in query_lower:
            return self.extract_instructions(chunks, combined_text)
        elif "deadline" in query_lower or "due date" in query_lower:
            return self.extract_deadline_info(chunks, combined_text)
        else:
            return self.provide_general_answer(query, chunks, combined_text)
            
    def extract_requirements(self, chunks: List[Dict], combined_text: str) -> str:
        """Extract requirements from the document context"""
        requirements = []
        lines = combined_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if any(keyword in line.lower() for keyword in ['requirement', 'must', 'should', 'need to', 'required']):
                if line and not line.startswith('#'):
                    requirements.append(line)
        
        if requirements:
            response = "Based on the documents, here are the requirements:\n\n"
            for i, req in enumerate(requirements[:10], 1):  # Limit to 10 requirements
                response += f"{i}. {req}\n"
            return response
        else:
            # Fallback to showing relevant sections
            relevant_sections = []
            for chunk in chunks:
                if any(keyword in chunk['text'].lower() for keyword in ['requirement', 'must', 'should', 'task']):
                    relevant_sections.append(f"From {chunk['filename']}:\n{chunk['text'][:300]}...")
            
            if relevant_sections:
                response = "Here are the relevant sections about requirements:\n\n"
                response += "\n\n".join(relevant_sections[:3])
                return response
            else:
                return "I found some relevant content, but couldn't identify specific requirements. Here's what I found:\n\n" + chunks[0]['text'][:500] + "..."
    
    def extract_citation_info(self, chunks: List[Dict], combined_text: str) -> str:
        """Extract citation information from the document context"""
        citation_info = []
        lines = combined_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if any(keyword in line.lower() for keyword in ['citation', 'cite', 'reference', 'bibliography', 'source']):
                if line and len(line) > 10:
                    citation_info.append(line)
        
        if citation_info:
            response = "Based on the documents, here's information about citations:\n\n"
            for info in citation_info[:5]:
                response += f"• {info}\n"
            return response
        else:
            return "I found some relevant content about citations:\n\n" + chunks[0]['text'][:500] + "..."
    
    def extract_definition(self, query: str, chunks: List[Dict], combined_text: str) -> str:
        """Extract definition or explanation for a term"""
        # Try to find the term being defined
        if "what is" in query.lower():
            term = query.lower().replace("what is", "").strip().rstrip("?")
        elif "define" in query.lower():
            term = query.lower().replace("define", "").strip()
        else:
            term = query.lower()
            
        # Look for definitions in the text
        lines = combined_text.split('.')
        definitions = []
        
        for line in lines:
            if term in line.lower() and any(keyword in line.lower() for keyword in ['is', 'means', 'refers to', 'defined as']):
                definitions.append(line.strip())
        
        if definitions:
            response = f"Based on the documents, here's what I found about '{term}':\n\n"
            for definition in definitions[:3]:
                response += f"• {definition}\n"
            return response
        else:
            return f"Here's the most relevant information I found about '{term}':\n\n" + chunks[0]['text'][:500] + "..."
    
    def extract_instructions(self, chunks: List[Dict], combined_text: str) -> str:
        """Extract how-to instructions or procedures"""
        instructions = []
        lines = combined_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if any(keyword in line.lower() for keyword in ['step', 'first', 'then', 'next', 'follow', 'procedure']):
                if line and len(line) > 5:
                    instructions.append(line)
        
        if instructions:
            response = "Here are the instructions I found:\n\n"
            for i, instruction in enumerate(instructions[:8], 1):
                response += f"{i}. {instruction}\n"
            return response
        else:
            return "Here's the relevant procedural information:\n\n" + chunks[0]['text'][:500] + "..."
    
    def extract_deadline_info(self, chunks: List[Dict], combined_text: str) -> str:
        """Extract deadline or date information"""
        deadline_info = []
        lines = combined_text.split('\n')
        
        for line in lines:
            line = line.strip()
            if any(keyword in line.lower() for keyword in ['deadline', 'due', 'date', 'submit', 'submission']):
                if line and len(line) > 10:
                    deadline_info.append(line)
        
        if deadline_info:
            response = "Here's the deadline information:\n\n"
            for info in deadline_info[:5]:
                response += f"• {info}\n"
            return response
        else:
            return "Here's the relevant information about deadlines:\n\n" + chunks[0]['text'][:500] + "..."
    
    def provide_general_answer(self, query: str, chunks: List[Dict], combined_text: str) -> str:
        """Provide a general answer when no specific pattern is matched"""
        # Try to find the most relevant chunk
        best_chunk = chunks[0]  # Default to first chunk
        
        # Look for chunks that contain query terms
        query_terms = query.lower().split()
        best_score = 0
        
        for chunk in chunks:
            score = sum(1 for term in query_terms if term in chunk['text'].lower())
            if score > best_score:
                best_score = score
                best_chunk = chunk
        
        response = f"Based on your question '{query}', here's the most relevant information I found:\n\n"
        response += best_chunk['text'][:800] + ("..." if len(best_chunk['text']) > 800 else "")
        response += f"\n\n(Source: {best_chunk['filename']})"
        
        return response

class CoordinatorAgent:
    """Main coordinator agent that orchestrates the workflow"""
    
    def __init__(self, message_bus: MCPMessageBus):
        self.message_bus = message_bus
        self.message_bus.subscribe("CoordinatorAgent", self.handle_message)
        self.active_traces = {}
        
    async def handle_message(self, message: MCPMessage):
        """Handle incoming MCP messages"""
        logger.info(f"CoordinatorAgent: Received {message.type} from {message.sender}")
        
        # Store message response in active_traces for UI retrieval
        self.active_traces[message.trace_id] = {
            "status": message.type, # Use message type as status for now
            "payload": message.payload
        }
        # You might want more granular status updates depending on your UI needs
        if message.type == "EMBEDDING_COMPLETE":
            st.session_state.indexed_document_names.add(message.payload.get("filename"))
            st.rerun() # Rerun to update the indexed documents list
        elif message.type == "FINAL_RESPONSE":
            st.session_state.chat_history.append({"role": "assistant", "content": message.payload["response"], "sources": message.payload["sources"]})
            st.session_state.response_ready = True # Signal that a response is ready
            st.rerun() # Rerun to display the response
        elif message.type == "ERROR":
            st.error(f"Error in trace {message.trace_id}: {message.payload.get('error')}")
            st.session_state.response_ready = True # Allow UI to unblock
            st.rerun() # Rerun to display error
        elif message.type == "INDEX_CLEARED":
            st.session_state.indexed_document_names.clear()
            st.session_state.chat_history = []
            st.success("All indexed data cleared!")
            st.rerun()
            
    async def process_document(self, uploaded_file: UploadedFile) -> str:
        """Process an uploaded document"""
        trace_id = str(uuid.uuid4())
        file_content = uploaded_file.read()
        
        message = MCPMessage(
            sender="CoordinatorAgent",
            receiver="IngestionAgent",
            type="INGEST_DOCUMENT",
            trace_id=trace_id,
            payload={
                "filename": uploaded_file.name,
                "type": uploaded_file.type,
                "content": file_content
            }
        )
        
        await self.message_bus.publish(message)
        return trace_id
    
    async def process_query(self, query: str) -> str:
        """Process a user query"""
        trace_id = str(uuid.uuid4())
        
        message = MCPMessage(
            sender="CoordinatorAgent",
            receiver="RetrievalAgent",
            type="RETRIEVE_CONTEXT",
            trace_id=trace_id,
            payload={"query": query, "top_k": 5} # Request top 5 chunks
        )
        
        await self.message_bus.publish(message)
        return trace_id

    async def clear_all_data(self) -> str:
        """Initiate clearing of all indexed data."""
        trace_id = str(uuid.uuid4())
        message = MCPMessage(
            sender="CoordinatorAgent",
            receiver="RetrievalAgent",
            type="CLEAR_INDEX",
            trace_id=trace_id,
            payload={}
        )
        await self.message_bus.publish(message)
        return trace_id

# --- Streamlit UI Class ---
class RAGSystemUI:
    def __init__(self):
        # Initialize message bus and agents, and store them in session_state
        # This prevents re-initialization on every rerun
        if "message_bus" not in st.session_state:
            st.session_state.message_bus = MCPMessageBus()
            st.session_state.ingestion_agent = IngestionAgent(st.session_state.message_bus)
            st.session_state.retrieval_agent = RetrievalAgent(st.session_state.message_bus)
            st.session_state.llm_response_agent = LLMResponseAgent(st.session_state.message_bus)
            st.session_state.coordinator_agent = CoordinatorAgent(st.session_state.message_bus)
            
            # Initialize chat history and indexed documents list
            st.session_state.chat_history = []
            st.session_state.indexed_document_names = set() # Use a set for unique names
            st.session_state.response_ready = True # Flag to control query input

        self.message_bus = st.session_state.message_bus
        self.coordinator_agent = st.session_state.coordinator_agent
        self.retrieval_agent = st.session_state.retrieval_agent # Used to check for indexed docs

    async def run(self):
        st.set_page_config(layout="wide", page_title="Intelligent RAG Chatbot")

        st.title("Intelligent RAG System with Agentic Workflow")

        # --- Document Upload Section ---
        st.header("Document Upload")
        uploaded_file = st.file_uploader(
            "Upload PDF, DOCX, PPTX, CSV, or TXT files",
            type=["pdf", "docx", "pptx", "csv", "txt", "md"],
            accept_multiple_files=False, # Process one at a time for simplicity
            key="file_uploader"
        )

        if uploaded_file:
            st.write(f"**Selected:** {uploaded_file.name} ({uploaded_file.size / 1024:.2f} KB)")
            
            # Check if this file is already in the indexed set
            if uploaded_file.name not in st.session_state.indexed_document_names:
                if st.button(f"Process {uploaded_file.name}", key="process_btn"):
                    with st.spinner(f"Processing '{uploaded_file.name}'... This may take a moment."):
                        try:
                            trace_id = await self.coordinator_agent.process_document(uploaded_file)
                            # Wait for the document processing to complete or error out
                            # This loop needs to be robust, possibly with a timeout
                            while trace_id not in self.coordinator_agent.active_traces or \
                                  self.coordinator_agent.active_traces[trace_id]["status"] not in ["embedding_complete", "ERROR"]:
                                await asyncio.sleep(0.1) # Short delay to prevent busy-waiting
                            
                            status_info = self.coordinator_agent.active_traces[trace_id]
                            if status_info["status"] == "embedding_complete":
                                st.success(f"'{uploaded_file.name}' processed and indexed successfully! Total chunks: {status_info['payload']['total_chunks_in_index']}")
                                st.session_state.indexed_document_names.add(uploaded_file.name) # Add to displayed list
                            else:
                                st.error(f"Failed to process '{uploaded_file.name}': {status_info['payload'].get('error', 'Unknown error')}")
                        except Exception as e:
                            st.error(f"An unexpected error occurred during document processing: {e}")
            else:
                st.info(f"'{uploaded_file.name}' is already indexed.")
        
        # Display Indexed Documents
        st.subheader("Indexed Documents")
        if self.retrieval_agent.index and self.retrieval_agent.index.ntotal > 0:
            st.write(f"Total chunks in index: **{self.retrieval_agent.index.ntotal}**")
            for doc_name in sorted(list(st.session_state.indexed_document_names)):
                st.write(f"- {doc_name}")
            
            if st.button("Clear All Indexed Data", key="clear_index_btn"):
                with st.spinner("Clearing all indexed documents..."):
                    trace_id = await self.coordinator_agent.clear_all_data()
                    while trace_id not in self.coordinator_agent.active_traces or \
                          self.coordinator_agent.active_traces[trace_id]["status"] != "INDEX_CLEARED":
                        await asyncio.sleep(0.1)
                    # The coordinator's handle_message already updates session state and reruns
        else:
            st.write("No documents indexed yet.")

        st.markdown("---")

        # --- Chat Interface Section ---
        st.header("Chat with your Documents")

        # Display chat messages
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
                if "sources" in message and message["sources"]:
                    st.caption(f"Sources: {', '.join(message['sources'])}")

        # Chat input
        user_query = st.chat_input(
            "Type your question here...",
            disabled=not st.session_state.response_ready # Disable input while processing
        )

        if user_query:
            # Add user query to chat history
            st.session_state.chat_history.append({"role": "user", "content": user_query})
            st.session_state.response_ready = False # Disable input
            
            # Display user message immediately
            with st.chat_message("user"):
                st.markdown(user_query)

            # Trigger RAG process
            with st.spinner("Searching and generating response..."):
                trace_id = await self.coordinator_agent.process_query(user_query)
                # Wait for the response to be ready
                while trace_id not in self.coordinator_agent.active_traces or \
                      self.coordinator_agent.active_traces[trace_id]["status"] not in ["FINAL_RESPONSE", "ERROR"]:
                    await asyncio.sleep(0.1) # Short delay to prevent busy-waiting
                
                # The CoordinatorAgent's handle_message for FINAL_RESPONSE/ERROR
                # will update chat_history and set response_ready=True, triggering a rerun.

        st.markdown("---")

        # --- System Log (MCP Messages) Section ---
        st.subheader("System Log (MCP Messages)")
        if self.message_bus.get_mcp_log():
            for msg in reversed(self.message_bus.get_mcp_log()): # Show latest first
                st.json(msg)
        else:
            st.info("No MCP messages yet.")

# --- Main application entry point ---
if __name__ == "__main__":
    # If GOOGLE_API_KEY is not set as an environment variable, prompt for it
    if not os.environ.get("GOOGLE_API_KEY") and LLM_AVAILABLE:
        api_key = st.sidebar.text_input("Enter your Google API Key:", type="password")
        if api_key:
            os.environ["GOOGLE_API_KEY"] = api_key
            st.sidebar.success("API Key set!")
            # Re-initialize LLM agent to pick up the new API key
            if "llm_response_agent" in st.session_state:
                st.session_state.llm_response_agent = LLMResponseAgent(st.session_state.message_bus)
                st.rerun()
        else:
            st.sidebar.warning("Please enter your Google API Key to enable LLM responses.")
            LLM_AVAILABLE = False # Temporarily disable if key not provided

    ui = RAGSystemUI()
    asyncio.run(ui.run())